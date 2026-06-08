"""文档构建服务 — 从上传文档中通过 LLM 多轮对话抽取本体"""
from __future__ import annotations

import io
import json
import logging
import re
import uuid
from typing import Generator

from openai import OpenAI

from app.config import settings
from app.services.ontology_constraints import build_constraint_prompt, validate_ontology_output, validate_and_retry

logger = logging.getLogger(__name__)

_sessions: dict[str, dict] = {}


def _get_llm_client(db=None) -> OpenAI:
    from app.services.llm_resolver import get_llm_client
    return get_llm_client(db=db, scene="ontology")


def _extract_json(text: str) -> str:
    text = re.sub(r"<think>[\s\S]*?</think>", "", text).strip()
    if text.startswith("```"):
        text = text.split("\n", 1)[-1].rsplit("```", 1)[0].strip()
    match = re.search(r"\{[\s\S]*\}", text)
    if match:
        return match.group(0)
    return text


def _sanitize_docx(data: bytes) -> bytes:
    """Remove broken relationship entries (e.g. Target='../NULL') from docx archive."""
    import zipfile
    from xml.etree import ElementTree as ET

    NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    src = zipfile.ZipFile(io.BytesIO(data))
    names = set(src.namelist())

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as dst:
        for item in src.infolist():
            raw = src.read(item.filename)
            if item.filename.endswith(".rels"):
                tree = ET.fromstring(raw)
                dirty = False
                for rel in list(tree):
                    target = rel.get("Target", "")
                    if target.startswith("http"):
                        continue
                    # resolve relative path against rels parent dir
                    import posixpath
                    rels_dir = posixpath.dirname(item.filename)
                    parent_dir = posixpath.dirname(rels_dir)
                    resolved = posixpath.normpath(posixpath.join(parent_dir, target))
                    if resolved not in names and target not in names:
                        tree.remove(rel)
                        dirty = True
                if dirty:
                    ET.register_namespace("", NS)
                    raw = ET.tostring(tree, encoding="unicode", xml_declaration=True).encode("utf-8")
            dst.writestr(item, raw)
    src.close()
    return buf.getvalue()


def parse_file_content(filename: str, content: bytes) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext in ("txt", "csv", "md"):
        return content.decode("utf-8", errors="ignore")[:8000]

    if ext == "pdf":
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                texts = [p.extract_text() or "" for p in pdf.pages[:20]]
            return "\n".join(texts)[:8000]
        except Exception as e:
            return f"[PDF 解析失败: {e}]"

    if ext in ("docx", "doc"):
        try:
            import docx
            doc = docx.Document(io.BytesIO(_sanitize_docx(content)))
            texts = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(texts)[:8000]
        except Exception as e:
            return f"[DOCX 解析失败: {e}]"

    if ext in ("xlsx", "xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets[:5]:
                lines.append(f"## Sheet: {ws.title}")
                for row in ws.iter_rows(max_row=100, values_only=True):
                    line = "\t".join(str(c) if c is not None else "" for c in row)
                    if line.strip():
                        lines.append(line)
            return "\n".join(lines)[:8000]
        except Exception as e:
            return f"[Excel 解析失败: {e}]"

    if ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides[:30]:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                texts.append(para.text)
            return "\n".join(texts)[:8000]
        except Exception as e:
            return f"[PPT 解析失败: {e}]"

    return f"[不支持的文件格式: {ext}]"


def create_session(files: list[dict]) -> dict:
    session_id = str(uuid.uuid4())
    _sessions[session_id] = {
        "id": session_id,
        "files": files,
        "history": [],
        "current_ontology": None,
    }
    return _sessions[session_id]


def get_session(session_id: str) -> dict | None:
    return _sessions.get(session_id)


_BASE_SYSTEM_PROMPT = """你是一名深耕中国联通业务体系的运营商领域专家兼本体建模专家，熟悉联通在公众客户、政企客户、网络资源、产品中心、订单中心、计费账务、客户服务、渠道运营、智慧家庭、物联网、云网融合、数据中台和智慧运营等领域的业务对象、流程规则和系统边界。你的任务是从用户提供的业务文档中提取本体（实体、属性、关系）。

规则：
1. 每次回复必须包含完整的本体 JSON
2. 在 JSON 之外可用自然语言解释你的分析思路
"""


def _build_system_prompt() -> str:
    return _BASE_SYSTEM_PROMPT + "\n" + build_constraint_prompt()


def chat_stream(
    session_id: str,
    message: str,
    business_desc: str,
) -> Generator[str, None, None]:
    session = _sessions.get(session_id)
    if not session:
        yield f"data: {json.dumps({'event': 'error', 'message': 'session not found'})}\n\n"
        return

    yield f"data: {json.dumps({'event': 'thinking', 'message': '正在分析文档内容...'})}\n\n"

    messages = [{"role": "system", "content": _build_system_prompt()}]

    if not session["history"]:
        doc_contents = "\n\n".join(
            f"### 文件: {f['name']}\n{f['content']}" for f in session["files"]
        )
        first_user_msg = f"""业务需求：{business_desc}

以下是上传的业务文档内容：

{doc_contents}

请从以上文档中提取本体的实体、属性和关系。"""
        messages.append({"role": "user", "content": first_user_msg})
        session["history"].append({"role": "user", "content": first_user_msg})
    else:
        for h in session["history"]:
            messages.append({"role": h["role"], "content": h["content"]})
        messages.append({"role": "user", "content": message})
        session["history"].append({"role": "user", "content": message})

    client = _get_llm_client()
    try:
        resp = client.chat.completions.create(
            model=settings.LLM_MODEL,
            messages=messages,
            temperature=0.3,
            stream=True,
        )
    except Exception as e:
        yield f"data: {json.dumps({'event': 'error', 'message': str(e)})}\n\n"
        return

    full_content = ""
    in_think = False
    buffer = ""
    think_notified = False

    for chunk in resp:
        delta = chunk.choices[0].delta
        if delta.content:
            full_content += delta.content
            buffer += delta.content

            while buffer:
                if in_think:
                    end_idx = buffer.find("</think>")
                    if end_idx == -1:
                        buffer = ""
                        break
                    else:
                        buffer = buffer[end_idx + 8:]
                        in_think = False
                else:
                    start_idx = buffer.find("<think>")
                    if start_idx == -1:
                        if "<" in buffer and not buffer.endswith(">"):
                            safe = buffer[:buffer.rfind("<")]
                            if safe:
                                yield f"data: {json.dumps({'event': 'token', 'content': safe})}\n\n"
                            buffer = buffer[len(safe):]
                        else:
                            yield f"data: {json.dumps({'event': 'token', 'content': buffer})}\n\n"
                            buffer = ""
                        break
                    else:
                        if start_idx > 0:
                            yield f"data: {json.dumps({'event': 'token', 'content': buffer[:start_idx]})}\n\n"
                        buffer = buffer[start_idx + 7:]
                        in_think = True
                        if not think_notified:
                            think_notified = True
                            yield f"data: {json.dumps({'event': 'thinking', 'message': 'AI 正在深度思考...'})}\n\n"

    session["history"].append({"role": "assistant", "content": full_content})

    ontology_json = _extract_json(full_content)
    if ontology_json:
        result, errors = validate_ontology_output(ontology_json)
        if result is not None:
            session["current_ontology"] = result
            yield f"data: {json.dumps({'event': 'ontology', 'data': result})}\n\n"
        else:
            yield f"data: {json.dumps({'event': 'validating', 'message': '正在校验修正输出...'})}\n\n"
            client = _get_llm_client()

            def retry_caller(retry_prompt: str) -> str:
                resp = client.chat.completions.create(
                    model=settings.LLM_MODEL,
                    messages=[
                        {"role": "system", "content": retry_prompt},
                        {"role": "user", "content": "请输出修正后的完整JSON。"},
                    ],
                    temperature=0,
                )
                return resp.choices[0].message.content or ""

            result, warnings = validate_and_retry(retry_caller, full_content, max_retries=5)
            if result is not None:
                session["current_ontology"] = result
                yield f"data: {json.dumps({'event': 'ontology', 'data': result})}\n\n"
                if warnings:
                    yield f"data: {json.dumps({'event': 'validation_warning', 'message': warnings})}\n\n"
            else:
                try:
                    fallback = json.loads(ontology_json)
                    session["current_ontology"] = fallback
                    yield f"data: {json.dumps({'event': 'ontology', 'data': fallback})}\n\n"
                    yield f"data: {json.dumps({'event': 'validation_warning', 'message': '输出未完全通过校验，请人工检查'})}\n\n"
                except json.JSONDecodeError:
                    pass

    yield "data: [DONE]\n\n"
